"""Filesystem helpers: walking directories, reading file contents, and
rendering directory trees."""

import os
import re
import xml.etree.ElementTree as ET
import zipfile

import docx
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation

from config import IMAGE_EXTENSIONS, TEXT_EXTENSIONS
from logging_setup import get_logger
from manifest import MANIFEST_NAME

logger = get_logger(__name__)


def read_text_file(file_path, max_chars=3000):
    """Read up to `max_chars` of a plain text file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read(max_chars)
    except OSError as exc:
        logger.error("Error reading text file %s: %s", file_path, exc)
        return None


def read_docx_file(file_path):
    """Read text content from a .docx or .doc file."""
    try:
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)
    except Exception as exc:
        logger.error("Error reading DOCX file %s: %s", file_path, exc)
        return None


def read_pdf_file(file_path, max_pages=3):
    """Read text content from the first `max_pages` pages of a PDF."""
    try:
        doc = fitz.open(file_path)
        pages = [doc.load_page(i).get_text() for i in range(min(max_pages, len(doc)))]
        return '\n'.join(pages)
    except Exception as exc:
        logger.error("Error reading PDF file %s: %s", file_path, exc)
        return None


def read_spreadsheet_file(file_path):
    """Read text content from an Excel or CSV file."""
    try:
        df = pd.read_csv(file_path) if file_path.lower().endswith('.csv') else pd.read_excel(file_path)
        return df.to_string()
    except Exception as exc:
        logger.error("Error reading spreadsheet file %s: %s", file_path, exc)
        return None


def read_ppt_file(file_path):
    """Read text content from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        full_text = [
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        ]
        return '\n'.join(full_text)
    except Exception as exc:
        logger.error("Error reading PowerPoint file %s: %s", file_path, exc)
        return None


_READERS_BY_EXTENSION = {
    '.txt': read_text_file,
    '.md': read_text_file,
    '.docx': read_docx_file,
    '.doc': read_docx_file,
    '.pdf': read_pdf_file,
    '.xls': read_spreadsheet_file,
    '.xlsx': read_spreadsheet_file,
    '.csv': read_spreadsheet_file,
    '.ppt': read_ppt_file,
    '.pptx': read_ppt_file,
}


def read_file_data(file_path):
    """Read content from a file based on its extension, or None if unsupported."""
    ext = os.path.splitext(file_path.lower())[1]
    reader = _READERS_BY_EXTENSION.get(ext)
    return reader(file_path) if reader else None


_CONTAINER_NS = {'c': 'urn:oasis:names:tc:opendocument:xmlns:container'}
_DC_NS = {'dc': 'http://purl.org/dc/elements/1.1/'}


def read_epub_metadata(file_path):
    """Extract {'title', 'creator', 'subjects'} from an epub's OPF metadata
    using only the stdlib (an epub is a zip with a Dublin Core manifest), or
    None if the file isn't a well-formed epub. Mobi/azw have no such cheap
    parse and are handled by filename alone."""
    try:
        with zipfile.ZipFile(file_path) as zf:
            container = ET.fromstring(zf.read('META-INF/container.xml'))
            rootfile = container.find('.//c:rootfile', _CONTAINER_NS)
            opf = ET.fromstring(zf.read(rootfile.get('full-path')))

        def text(tag):
            el = opf.find(f'.//dc:{tag}', _DC_NS)
            return el.text.strip() if el is not None and el.text else None

        return {
            'title': text('title'),
            'creator': text('creator'),
            'subjects': [
                el.text.strip() for el in opf.findall('.//dc:subject', _DC_NS)
                if el.text and el.text.strip()
            ],
        }
    except Exception as exc:
        logger.debug("Could not read epub metadata from %s: %s", file_path, exc)
        return None


def display_directory_tree(path):
    """Print a directory tree similar to the `tree` command, prefixed by the full path."""
    def tree(dir_path, prefix=''):
        contents = sorted(c for c in os.listdir(dir_path) if not c.startswith('.'))
        pointers = ['├── '] * (len(contents) - 1) + ['└── '] if contents else []
        for pointer, name in zip(pointers, contents):
            full_path = os.path.join(dir_path, name)
            print(prefix + pointer + name)
            if os.path.isdir(full_path):
                extension = '│   ' if pointer == '├── ' else '    '
                tree(full_path, prefix + extension)

    print(os.path.abspath(path))
    if os.path.isdir(path):
        tree(path)


def collect_file_paths(base_path):
    """Collect all file paths under `base_path` (or itself, if it's a file),
    excluding hidden files. Subtrees containing a manifest marker -- outputs
    of a previous run of this tool -- are skipped entirely, so an organized
    tree sitting inside the input never gets re-organized (which would feed
    generated category folders back into classification)."""
    if os.path.isfile(base_path):
        return [base_path]

    file_paths = []
    for root, dirs, files in os.walk(base_path):
        if MANIFEST_NAME in files:
            if root == os.path.normpath(base_path) or os.path.samefile(root, base_path):
                logger.warning(
                    "Input directory %s is itself a previous output of this tool; "
                    "re-organizing it anyway since it was chosen explicitly", base_path,
                )
            else:
                logger.warning("Skipping previously organized subtree: %s", root)
                dirs[:] = []
                continue
        file_paths.extend(os.path.join(root, name) for name in files if not name.startswith('.'))
    return file_paths


def separate_files_by_type(file_paths):
    """Split file paths into (image_files, text_files, other_files) based on
    extension. `other_files` holds extensions content mode can't read (e.g.
    archives, ebooks) -- the caller is responsible for not dropping them."""
    image_files, text_files, other_files = [], [], []
    for fp in file_paths:
        ext = os.path.splitext(fp.lower())[1]
        if ext in IMAGE_EXTENSIONS:
            image_files.append(fp)
        elif ext in TEXT_EXTENSIONS:
            text_files.append(fp)
        else:
            other_files.append(fp)
    return image_files, text_files, other_files


# Generic/auto-generated folder names that carry no real information about
# their contents (OS defaults, camera/app dumps, bare dates or numbers), so a
# folder with one of these names should NOT be trusted as a meaningful label.
_JUNK_FOLDER_NAMES = {
    'new folder', 'untitled folder', 'untitled', 'temp', 'tmp', 'desktop',
    'downloads', 'download', 'documents', 'my documents', 'pictures', 'my pictures',
    'screenshots', 'screenshot', 'misc', 'miscellaneous', 'stuff', 'dcim',
    'camera uploads', 'photos', 'files', 'backup', 'old', 'others', 'other',
    'unsorted', 'unnamed', 'inbox', 'attachments', 'exports', 'export',
}

_JUNK_FOLDER_PATTERNS = [re.compile(p) for p in (
    r'^new folder(\s*\(\d+\))?$',
    r'^untitled(\s*folder)?\s*\d*$',
    r'^copy of .*$',
    r'^img[-_ ]?\d+$',                # IMG_1234, IMG-1234
    r'^dcim\d*$',
    r'^\d{3,4}[a-z]+$',                # 100ANDRO, 101APPLE style camera dirs
    r'^screenshots?([-_ ]?\d+)?$',
    r'^\d{4}[-_]\d{2}[-_]\d{2}$',       # date-stamped, e.g. 2024-01-15
    r'^\d{8}$',                        # 20240115
    r'^\(?\d+\)?$',                    # purely numeric, e.g. "1", "(2)"
)]


def is_junk_folder_name(name):
    """Return True if `name` looks like a generic/auto-generated folder name
    rather than a meaningful, user-chosen label (e.g. "Downloads", "IMG_1234",
    "New Folder (2)") that shouldn't be trusted to describe its contents."""
    normalized = name.strip().lower()
    if normalized in _JUNK_FOLDER_NAMES:
        return True
    return any(pattern.match(normalized) for pattern in _JUNK_FOLDER_PATTERNS)


def split_by_folder_trust(file_paths, base_path):
    """Split files into (trusted, untrusted) based on whether the folder(s)
    they already sit in look like meaningful, user-chosen labels.

    A file is "trusted" (its existing location is kept as-is, skipping AI
    classification) only if it sits inside at least one subfolder of
    `base_path` and none of the folder names between `base_path` and the file
    look like junk. Files sitting directly in `base_path`, or under any junk
    folder name, are "untrusted" and go through the normal content pipeline.
    """
    trusted, untrusted = [], []
    for fp in file_paths:
        rel_dir = os.path.relpath(os.path.dirname(fp), base_path)
        if rel_dir == '.':
            untrusted.append(fp)
            continue
        components = rel_dir.split(os.sep)
        if any(is_junk_folder_name(c) for c in components):
            untrusted.append(fp)
        else:
            trusted.append(fp)
    return trusted, untrusted
